# backend/extract_text.py
import sys, os, boto3, io, json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
from dotenv import load_dotenv

load_dotenv()

# Set Tesseract CMD path - modify this path to match your installation
if os.name == 'nt':  # Windows
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def check_tesseract():
    """Verify Tesseract installation and configuration."""
    try:
        # Try a simple OCR operation to verify Tesseract works
        test_image = Image.new('RGB', (100, 30), color='white')
        pytesseract.image_to_string(test_image)
        return True
    except Exception as e:
        print(json.dumps({"error": f"Tesseract not properly configured: {str(e)}. Please ensure Tesseract is installed."}), file=sys.stderr)
        return False

def perform_ocr(image_bytes):
    """Performs OCR on image bytes and returns extracted text."""
    try:
        image = Image.open(io.BytesIO(image_bytes))
        
        # Convert RGBA to RGB if needed (Tesseract prefers RGB)
        if image.mode == 'RGBA':
            background = Image.new('RGB', image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[3])
            image = background
        
        # Attempt to improve image quality for OCR
        image = image.convert('L')  # Convert to grayscale
        
        # Run OCR with specific configuration for better results
        custom_config = r'--oem 3 --psm 6'  # Page segmentation mode 6: Assume uniform block of text
        text = pytesseract.image_to_string(image, config=custom_config)
        extracted = text.strip()
        
        if extracted:
            print(f"OCR succeeded, extracted {len(extracted)} characters", file=sys.stderr)
            return extracted
        else:
            print("OCR produced no text", file=sys.stderr)
            return ""
    except Exception as e:
        print(f"OCR Error: {e}", file=sys.stderr)
        return ""

def download_from_s3(bucket_name, s3_key):
    try:
        s3 = boto3.client('s3')
        buf = io.BytesIO()
        s3.download_fileobj(bucket_name, s3_key, buf)
        buf.seek(0)
        return buf
    except Exception as e:
        print(json.dumps({"error": f"S3 download failed: {e}"}), file=sys.stderr)
        sys.exit(1)

def extract_text_from_presentation(file_buffer):
    data = []
    try:
        # First verify Tesseract is working
        if not check_tesseract():
            return data

        prs = Presentation(file_buffer)
        if not prs.slides: return data

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            slide_texts = []
            notes_text = None
            image_count = 0
            ocr_count = 0

            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                    try:
                        image_bytes = shape.image.blob
                        ocr_text = perform_ocr(image_bytes)
                        if ocr_text:
                            ocr_count += 1
                            slide_texts.append(f"[OCR Text: {ocr_text}]")
                            print(f"Slide {slide_num}: OCR successful for image {image_count}", file=sys.stderr)
                    except Exception as img_err:
                        print(f"Error processing image {image_count} on slide {slide_num}: {img_err}", file=sys.stderr)

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            # Include OCR statistics in the slide data
            slide_data = {
                "slide": slide_num,
                "text": " ".join(slide_texts),
                "notes": notes_text,
                "stats": {
                    "total_images": image_count,
                    "ocr_successful": ocr_count
                }
            }
            data.append(slide_data)
            
            # Log slide processing summary
            print(f"Processed slide {slide_num}: {image_count} images, {ocr_count} OCR successes", file=sys.stderr)

        return data
    except Exception as e:
        print(json.dumps({"error": f"PPTX processing failed: {e}"}), file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(json.dumps({"error": "Usage: python extract_text.py <s3_key>"}), file=sys.stderr)
        sys.exit(1)

    s3_key = sys.argv[1]
    bucket_name = os.getenv("S3_BUCKET_NAME")
    if not bucket_name:
        print(json.dumps({"error": "S3_BUCKET_NAME environment variable not set"}), file=sys.stderr)
        sys.exit(1)

    presentation_buffer = download_from_s3(bucket_name, s3_key)
    extracted_content = extract_text_from_presentation(presentation_buffer)
    print(json.dumps({"data": extracted_content}))